import os
import sys
sys.stdout.reconfigure(encoding='utf-8')

try:
    from pptx import Presentation
except ImportError:
    import subprocess
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation

def inspect_pptx(path):
    print(f"\n==========================================")
    print(f"Reading PPTX: {path}")
    try:
        prs = Presentation(path)
        for idx, slide in enumerate(prs.slides):
            print(f"\n--- SLIDE {idx+1} ---")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            print(f"  {text}")
    except Exception as e:
        print(f"Error: {e}")

def main():
    folder = r"C:\Users\giack\Desktop\Progetto TedX"
    files = [
        "MyTEDx_CareerPath_AI_Light (1).pptx",
        "MyTEDx_CareerPathAI_Compito1.pptx",
        "MyTEDx_CareerPathAI_Compito2.pptx"
    ]
    for file in files:
        path = os.path.join(folder, file)
        if os.path.exists(path):
            inspect_pptx(path)

if __name__ == '__main__':
    main()
