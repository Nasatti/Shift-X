import os
import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation

def main():
    pptx_path = r"C:\Users\giack\Desktop\Progetto TedX\MyTEDx_CareerPath_AI_Light (1).pptx"
    out_path = r"C:\Users\giack\.gemini\antigravity\scratch\shift_x_flutter\MyTEDx_CareerPath_AI_Light.txt"
    
    if not os.path.exists(pptx_path):
        print(f"File not found: {pptx_path}")
        return
        
    print(f"Reading {pptx_path} and writing to {out_path}...")
    try:
        prs = Presentation(pptx_path)
        with open(out_path, 'w', encoding='utf-8') as f:
            for idx, slide in enumerate(prs.slides):
                f.write(f"\n--- SLIDE {idx+1} ---\n")
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text = paragraph.text.strip()
                            if text:
                                f.write(f"{text}\n")
        print("Done!")
    except Exception as e:
        print(f"Error: {e}")

if __name__ == '__main__':
    main()
